#!/usr/bin/env python3
"""SNC Environmental Consultancy - Corporate deck built to the official brand
guidelines (Kamerik 105 / Golos Text, Deep Vision Navy / Innovation Blue /
Sustainability Green / Energy Orange / Pure Harmony Beige) and the SNC
presentation template layout."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ------------------------------------------------ brand system (from guidelines)
NAVY   = RGBColor(0x2F, 0x48, 0x58)   # Deep Vision Navy  #2F4858 (primary)
BLUE   = RGBColor(0x00, 0x6A, 0x91)   # Innovation Blue   #006A91 (secondary)
GREEN  = RGBColor(0x8A, 0xB9, 0x25)   # Sustainability Grn #8AB925 (accent)
ORANGE = RGBColor(0xF8, 0x5E, 0x00)   # Energy Orange    #F85E00 (accent)
BEIGE  = RGBColor(0xEF, 0xED, 0xE8)   # Pure Harmony Beige #EFEDE8 (neutral)
INK    = RGBColor(0x3A, 0x4A, 0x55)   # body text (harmonised)
MUTED  = RGBColor(0x6B, 0x76, 0x7B)   # secondary text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NEARWH = RGBColor(0xE9, 0xEC, 0xE9)   # card fill on beige
CARD   = RGBColor(0xF6, 0xF4, 0xF0)   # light card
LINEC  = RGBColor(0xD9, 0xD6, 0xCF)   # hairline on beige

# Brand fonts (Kamerik 105 / Golos Text as per guidelines; fall back gracefully)
HEAD = "Kamerik 105"
BODY = "Golos Text"

SW, SH = Inches(13.333), Inches(7.5)
ML     = Inches(0.0)
MR     = Inches(0.0)
REG_LINE = "NCEC Cat. A License 123456  \u00b7  CR 2050 123 456  \u00b7  VAT 3107 234 567 81  \u00b7  Aligned with MEWA Environmental Law"

LOCK_NAVY  = 'assets/lockup_navy.png'
LOCK_GREEN = 'assets/lockup_green.png'
RING       = 'assets/brand_ring.png'
TEXT_LIGHT = 'assets/texture_light.png'


def font(color, size, bold=False, head=False, italic=False, spacing=None):
    def apply(run):
        run.font.name = HEAD if head else BODY
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if spacing is not None:
            run.font._rPr.set('spc', str(int(spacing * 100)))
    return apply


def add_text(slide, x, y, w, h, paragraphs, anchor=MSO_ANCHOR.TOP,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    for i, p in enumerate(paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', align)
        para.space_before = Pt(p.get('space_before', 0))
        para.space_after = Pt(p.get('space_after', 0))
        if p.get('line'):
            para.line_spacing = p['line']
        for run_spec in p['runs']:
            r = para.add_run()
            apply = font(run_spec[2], run_spec[1], bold=run_spec[3],
                         head=run_spec[4] if len(run_spec) > 4 else False)
            apply(r)
            r.text = run_spec[0]
            if len(run_spec) > 5 and run_spec[5]:
                r.font.italic = True
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
             shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w if line_w else 1)
    sp.shadow.inherit = False
    return sp


def pic(slide, path, x, y, w=None, h=None):
    if w is not None and h is not None:
        return slide.shapes.add_picture(path, x, y, width=w, height=h)
    if w is not None:
        return slide.shapes.add_picture(path, x, y, width=w)
    return slide.shapes.add_picture(path, x, y, height=h)


def bg(slide, color=BEIGE):
    add_rect(slide, 0, 0, SW, SH, color)


def header(slide, title, subtitle=None, title_size=40):
    """Centered heading, letter-spaced label style, template look."""
    add_logo_topright(slide)
    if subtitle:
        add_text(slide, Inches(0.9), Inches(1.45), SW - Inches(1.8), Inches(0.5),
                 [{'runs': [(subtitle.upper(), 12, MUTED, True, False)], 'align': PP_ALIGN.CENTER,
                   'space_after': 0}])
        ty = Inches(1.62)
    else:
        ty = Inches(1.55)
    add_text(slide, Inches(0.7), ty, SW - Inches(1.4), Inches(1.2),
             [{'runs': [(title, title_size, NAVY, True, True)], 'align': PP_ALIGN.CENTER, 'line': 1.05}])
    return ty + Inches(0.0)


def add_logo_topright(slide, w=Inches(1.55)):
    pic(slide, LOCK_NAVY, SW - Inches(2.6), Inches(0.42), w=w)


def footer(slide, page_no, dark=False):
    if dark:
        txt_col, ln = NEARWH, RGBColor(0x4A, 0x5C, 0x68)
    else:
        txt_col, ln = MUTED, LINEC
    add_rect(slide, Inches(0.55), Inches(6.9), SW - Inches(1.1), Emu(9525), ln)
    pic(slide, RING, Inches(0.55), Inches(7.02), w=Inches(0.2))
    add_text(slide, Inches(0.82), Inches(7.04), Inches(4.4), Inches(0.3),
             [{'runs': [("SNC Environmental Consultancy", 9, MUTED if not dark else NEARWH, True)]}])
    add_text(slide, Inches(0.55), Inches(7.04), SW - Inches(1.7), Inches(0.3),
             [{'runs': [(REG_LINE, 7.5, MUTED if not dark else NEARWH, False)],
               'align': PP_ALIGN.RIGHT}])
    add_text(slide, SW - Inches(1.2), Inches(7.04), Inches(0.5), Inches(0.3),
             [{'runs': [(str(page_no), 9.5, BLUE if not dark else GREEN, True)],
               'align': PP_ALIGN.RIGHT}])


def diagonal_watermark(slide, x, y, w, show_full=True):
    pic(slide, 'assets/ring_faint.png', x, y, w=w)


def photo_panel(slide, img, x, y, w, h, caption, accent=NAVY, sub=None):
    """Professionally embedded photograph: photo + navy caption bar + accent."""
    # image
    pic(slide, img, x, y, w=w, h=h)
    # thin frame
    add_rect(slide, x, y, w, h, None, line=RGBColor(0xC9, 0xC6, 0xBF), line_w=1)
    # navy caption bar anchored to the bottom of the frame
    bar_h = Inches(0.32)
    add_rect(slide, x, y + h - bar_h, w, bar_h, NAVY)
    add_text(slide, x + Inches(0.16), y + h - bar_h, w - Inches(0.32), bar_h,
             [{'runs': [(caption, 9.5, WHITE, True, False)]}], anchor=MSO_ANCHOR.MIDDLE)
    # brand accent underline
    add_rect(slide, x, y, Inches(0.07), h, accent)


prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
LAY = prs.slide_layouts[6]


def new_slide(bgcolor=BEIGE):
    s = prs.slides.add_slide(LAY)
    bg(s, bgcolor)
    return s


# ============================================================ 1. COVER
s = new_slide()
# LEFT navy panel (about 33%)
add_rect(s, 0, 0, Inches(4.6), SH, NAVY)
pic(s, LOCK_GREEN, Inches(0.6), Inches(0.6), w=Inches(2.6))
add_text(s, Inches(0.6), Inches(2.9), Inches(3.6), Inches(2.2),
         [{'runs': [("Company", 44, WHITE, True, True)], 'line': 1.02},
          {'runs': [("Profile", 44, WHITE, True, True)], 'line': 1.02}])
add_text(s, Inches(0.6), Inches(4.6), Inches(3.7), Inches(1.2),
         [{'runs': [("ENVIRONMENTAL CONSULTANCY", 12.5, GREEN, True, False)], 'space_after': 10},
          {'runs': [("Saudi-based Category A Environmental Consultancy", 13, NEARWH, False)],
           'line': 1.2}])
add_rect(s, Inches(0.6), Inches(6.4), Inches(2.2), Inches(0.04), ORANGE)
add_text(s, Inches(0.6), Inches(6.6), Inches(4.0), Inches(0.5),
         [{'runs': [("www.sncconsultancy.com", 13, WHITE, True)], 'line': 1.1}])
# RIGHT beige area with diagonal texture
pic(s, TEXT_LIGHT, Inches(4.6), 0, w=SW - Inches(4.6), h=SH)
# faint ring watermark bottom-right corner (away from text)
pic(s, 'assets/ring_faint.png', Inches(10.6), Inches(4.7), w=Inches(3.2))
add_text(s, Inches(5.0), Inches(2.2), Inches(4.8), Inches(2.0),
         [{'runs': [("Practical Environmental Solutions.", 21, NAVY, True, False)], 'line': 1.32},
          {'runs': [("Regulatory Clarity.  Project Readiness.", 21, NAVY, True, False)], 'line': 1.32}])
add_rect(s, Inches(5.0), Inches(4.05), Inches(3.4), Inches(0.05), ORANGE)
add_text(s, Inches(5.0), Inches(4.3), Inches(5.0), Inches(1.0),
         [{'runs': [("Environmental permitting \u00b7 EIA & studies \u00b7 EMP / CEMP \u00b7 "
                     "Compliance & reporting \u00b7 Waste \u00b7 Monitoring", 12.5, INK, False)], 'line': 1.3}])
add_text(s, Inches(5.0), Inches(6.35), Inches(7.6), Inches(0.4),
         [{'runs': [("Dammam, Eastern Province, KSA    \u00b7    +966 55 477 7412    \u00b7    info@snc.com.sa",
                     12, MUTED, False)]}])

# ============================================================ 2. WELCOME (template p2)
# Balanced two-column grid: shared top margin, matched column heights,
# accent bars aligned to the same left margin, no text crowding the photo.
s = new_slide()
diagonal_watermark(s, Inches(-0.6), Inches(-0.4), Inches(4.2))
pic(s, LOCK_NAVY, SW - Inches(3.0), Inches(0.42), w=Inches(2.0))

LX, LW = Inches(0.95), Inches(5.95)      # left text column
RX, RW = Inches(7.85), Inches(4.85)      # right photo column
BAR = Inches(1.9)                          # consistent accent-bar width

# RIGHT column first (defines the visual anchor): photo + caption block
photo_panel(s, 'assets/photo_dammam.png', RX, Inches(1.75),
            RW, Inches(2.9), "Dammam, Eastern Province \u2014 our home base", accent=BLUE)
add_rect(s, RX, Inches(4.95), Inches(0.9), Inches(0.05), GREEN)
add_text(s, RX, Inches(5.12), RW, Inches(0.6),
         [{'runs': [("ENVIRONMENTAL CONSULTANCY", 11.5, GREEN, True, False)], 'line': 1.1},
          {'runs': [("Licensed \u00b7 Registered \u00b7 MEWA-aligned", 10.5, MUTED, False)],
           'space_before': 4, 'line': 1.1}])

# LEFT column: heading + accent bar + intro + MEWA tag
add_text(s, LX, Inches(1.85), LW, Inches(0.7),
         [{'runs': [("Welcome", 42, NAVY, True, True)], 'line': 1.0}])
add_rect(s, LX, Inches(2.62), BAR, Inches(0.05), ORANGE)
add_text(s, LX, Inches(2.88), LW, Inches(1.9),
         [{'runs': [("SNC Environmental Consultancy is a Saudi-based Category A environmental consultancy "
                     "headquartered in Dammam, Eastern Province. From a single project to a national portfolio, "
                     "we help project owners, developers, contractors and operating facilities obtain and "
                     "maintain environmental approvals with confidence \u2014 translating complex NCEC and MEWA "
                     "requirements into a clear, documented, approval-ready plan.", 15, INK, False)], 'line': 1.4}])
add_rect(s, LX, Inches(4.95), BAR, Inches(0.05), GREEN)
add_text(s, LX, Inches(5.12), LW, Inches(0.6),
         [{'runs': [("ALIGNED WITH MEWA ENVIRONMENTAL LAW", 12, GREEN, True, False)], 'line': 1.1}])
footer(s, 2)

# ============================================================ 3. ABOUT US (template p3)
s = new_slide()
y = header(s, "About SNC", "Who we support")
diagonal_watermark(s, Inches(-0.5), Inches(4.0), Inches(3.4))
gw = (SW - Inches(2.4) - Inches(0.5)) / 3
c0 = Inches(0.9)
cy = Inches(2.55)
cards = [
    ("Project Owners & Developers", "Supporting projects from early environmental requirements through permitting and documentation.", BLUE),
    ("Operating Facilities", "Supporting active facilities with environmental records, compliance documentation and regulatory coordination.", GREEN),
    ("Public & Semi-Government Entities", "Environmental consultancy support for public, semi-government and strategic development projects.", ORANGE),
]
for i, (t, d, col) in enumerate(cards):
    x = c0 + i * (gw + Inches(0.25))
    add_rect(s, x, cy, gw, Inches(2.5), CARD, line=LINEC, line_w=1)
    add_rect(s, x, cy, gw, Inches(0.09), col)
    add_rect(s, x + Inches(0.32), cy + Inches(0.4), Inches(0.92), Inches(0.36), BEIGE)
    add_text(s, x + Inches(0.32), cy + Inches(0.4), Inches(0.92), Inches(0.36),
             [{'runs': [("SNC", 12, col, True, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.32), cy + Inches(0.92), gw - Inches(0.64), Inches(0.9),
             [{'runs': [(t, 15.5, NAVY, True, True)], 'line': 1.05}])
    add_text(s, x + Inches(0.32), cy + Inches(1.78), gw - Inches(0.64), Inches(0.7),
             [{'runs': [(d, 12, INK, False)], 'line': 1.15}])
add_rect(s, c0, Inches(5.35), SW - Inches(1.8), Inches(0.85), NAVY)
add_text(s, c0 + Inches(0.3), Inches(5.35), SW - Inches(2.4), Inches(0.85),
         [{'runs': [("Why SNC  ", 13, GREEN, True, False),
                    ("we are a Saudi-licensed Category A firm that turns environmental complexity into "
                     "approval-ready documentation \u2014 combining hands-on authority experience, multi-sector "
                     "delivery and full-lifecycle support with accredited, auditable laboratory coordination.",
                     13, WHITE, False)], 'line': 1.15}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3)

# ============================================================ 4. MISSION (template p4)
# Consistent two-column layout shared with VISION (p5) so they align as a pair.
s = new_slide()
y = header(s, "Our Mission", "Guided by Vision 2030")
diagonal_watermark(s, Inches(-0.6), Inches(-0.5), Inches(3.4))
LX, LW = Inches(0.9), Inches(5.5)          # left column (statement)
RX, RW = Inches(6.9), Inches(5.5)          # right column (cards)
# Left: mission statement, vertically centred in the content band
add_text(s, LX, Inches(2.55), LW, Inches(2.5),
         [{'runs': [("To provide innovative, data-driven and sustainable environmental solutions that "
                     "support a greener, more resilient future for industries, governments and communities "
                     "worldwide.", 20, NAVY, True, True)], 'line': 1.3}])
add_rect(s, LX, Inches(5.05), Inches(2.0), Inches(0.05), ORANGE)
add_text(s, LX, Inches(5.25), LW, Inches(1.4),
         [{'runs': [("We integrate economic growth with ecological resilience \u2014 fostering accountable, "
                     "sustainable and shared prosperity for every client and community we serve.",
                     12.5, INK, False)], 'line': 1.4}])
# Right: two stacked cards, evenly filling the same band
items = [
    ("Innovative Solutions", "Data-driven environmental solutions that are practical and implementable.", GREEN, "01"),
    ("Sustainability First", "Responsible progress that unites strategic innovation with environmental stewardship.", BLUE, "02"),
]
cy = Inches(2.5)
card_h, gap = Inches(1.9), Inches(0.28)
for i, (t, d, col, num) in enumerate(items):
    yy = cy + i * (card_h + gap)
    add_rect(s, RX, yy, RW, card_h, CARD, line=LINEC, line_w=1)
    add_rect(s, RX, yy, Inches(0.09), card_h, col)
    add_rect(s, RX + Inches(0.32), yy + Inches(0.36), Inches(0.52), Inches(0.52), col)
    add_text(s, RX + Inches(0.32), yy + Inches(0.36), Inches(0.52), Inches(0.52),
             [{'runs': [(num, 14, WHITE, True, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, RX + Inches(1.02), yy + Inches(0.4), RW - Inches(1.4), Inches(0.5),
             [{'runs': [(t, 15.5, NAVY, True, True)], 'line': 1.02}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, RX + Inches(0.32), yy + Inches(1.06), RW - Inches(0.64), Inches(0.75),
             [{'runs': [(d, 11.5, INK, False)], 'line': 1.15}])
footer(s, 4)

# ============================================================ 5. VISION (template p5)
# Mirrors MISSION (p4): same two-column geometry for a consistent pair.
s = new_slide()
y = header(s, "Our Vision", "A greener, cleaner future")
diagonal_watermark(s, Inches(-0.6), Inches(-0.5), Inches(3.4))
LX, LW = Inches(0.9), Inches(5.5)
RX, RW = Inches(6.9), Inches(5.5)
# Left: vision statement
add_text(s, LX, Inches(2.55), LW, Inches(2.5),
         [{'runs': [("We strive to drive responsible progress and become a leader in clean energy "
                     "solutions \u2014 with an international presence and a reputation for excellence.",
                     20, NAVY, True, True)], 'line': 1.3}])
add_rect(s, LX, Inches(4.42), Inches(2.0), Inches(0.05), ORANGE)
add_text(s, LX, Inches(4.54), LW, Inches(0.9),
         [{'runs': [("Through strong partnerships and evidence-based solutions we grow an economy that is "
                     "accountable, sustainable and built to last.", 12.5, INK, False)],
           'line': 1.32}])
# Professional photo panel: AlUla natural environment
photo_panel(s, 'assets/photo_alula.png', LX, Inches(5.34),
            Inches(5.5), Inches(1.35), "AlUla \u2014 the natural environment we help protect", accent=GREEN)
# Right: three value cards (mirrors the mission's numbered cards)
vals = [("Sustainability", "Environmental responsibility for a greener future.", GREEN, "01"),
        ("Innovation", "Adaptability and forward-thinking clean energy solutions.", BLUE, "02"),
        ("Integrity", "Collaboration, transparency, excellence and accountability.", NAVY, "03")]
cy = Inches(2.5)
card_h, gap = Inches(1.18), Inches(0.2)
for i, (t, d, col, num) in enumerate(vals):
    yy = cy + i * (card_h + gap)
    add_rect(s, RX, yy, RW, card_h, CARD, line=LINEC, line_w=1)
    add_rect(s, RX, yy, Inches(0.09), card_h, col)
    add_rect(s, RX + Inches(0.32), yy + Inches(0.32), Inches(0.44), Inches(0.44), col)
    add_text(s, RX + Inches(0.32), yy + Inches(0.32), Inches(0.44), Inches(0.44),
             [{'runs': [(num, 12.5, WHITE, True, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, RX + Inches(0.9), yy + Inches(0.24), RW - Inches(1.25), Inches(0.4),
             [{'runs': [(t, 14, NAVY, True, True)], 'line': 1.0}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, RX + Inches(0.9), yy + Inches(0.58), RW - Inches(1.25), Inches(0.55),
             [{'runs': [(d, 10.5, INK, False)], 'line': 1.1}])
footer(s, 5)

# ============================================================ 6. SERVICES (template p6)
s = new_slide()
y = header(s, "Our Services", "The full project lifecycle in one place", title_size=38)
svcs = [
    ("Environmental Permitting", "Secure construction and operational permits through managed submissions, authority responses and follow-up.", BLUE),
    ("EIA & Studies", "Environmental Impact Assessments and studies that identify impacts and plan practical mitigation.", GREEN),
    ("EMP / CEMP / OEMP", "Management plans for construction and operation with monitoring and control measures built in.", ORANGE),
    ("Compliance & Reporting", "Environmental records and periodic reporting that keep your facility audit-ready.", BLUE),
    ("Waste Advisory", "Waste management plans and classification guidance aligned to regulatory requirements.", GREEN),
    ("Monitoring", "Air, noise, water, soil and sediment monitoring coordinated through accredited laboratories.", ORANGE),
]
gw = (SW - Inches(2.6) - Inches(0.5)) / 3
c0 = Inches(0.85)
cy = Inches(2.45)
gh = Inches(1.9)
for i, (t, d, col) in enumerate(svcs):
    r, c = divmod(i, 3)
    x = c0 + c * (gw + Inches(0.25))
    yy = cy + r * (gh + Inches(0.28))
    add_rect(s, x, yy, gw, gh, CARD, line=LINEC, line_w=1)
    add_rect(s, x, yy, gw, Inches(0.09), col)
    add_rect(s, x + Inches(0.3), yy + Inches(0.34), Inches(0.86), Inches(0.34), BEIGE)
    add_text(s, x + Inches(0.3), yy + Inches(0.34), Inches(0.86), Inches(0.34),
             [{'runs': [("SNC", 11, col, True, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.3), yy + Inches(0.82), gw - Inches(0.6), Inches(0.6),
             [{'runs': [(t, 15.5, NAVY, True, True)], 'line': 1.04}])
    add_text(s, x + Inches(0.3), yy + Inches(1.42), gw - Inches(0.6), Inches(0.5),
             [{'runs': [(d, 11, INK, False)], 'line': 1.15}])
footer(s, 6)

# ============================================================ 7. REGULATORY APPROACH
s = new_slide()
y = header(s, "Our Regulatory Approach", "Structured from screening to compliance")
steps = [
    ("Regulatory Screening", "Early review of project activities, site context and expected permitting requirements.", BLUE, "01"),
    ("Documentation Readiness", "Preparation of environmental studies, plans, records and supporting documents.", GREEN, "02"),
    ("Authority Coordination", "Support during submission, follow-up, clarification and response to comments.", ORANGE, "03"),
    ("Compliance Continuity", "Post-permit support through records, reporting, monitoring and compliance documentation.", BLUE, "04"),
]
gw = (SW - Inches(2.5) - Inches(0.4)) / 2
c0 = Inches(0.9)
# supporting line under the title
add_text(s, Inches(1.0), Inches(2.24), SW - Inches(2.0), Inches(0.3),
         [{'runs': [("A documented, auditable path from first screening through to post-permit compliance \u2014 "
                     "so you maintain control at every stage.", 12.5, INK, False)], 'align': PP_ALIGN.CENTER}])
cy = Inches(2.7)
gh = Inches(1.85)
for i, (t, d, col, num) in enumerate(steps):
    r, c = divmod(i, 2)
    x = c0 + c * (gw + Inches(0.4))
    yy = cy + r * (gh + Inches(0.3))
    add_rect(s, x, yy, gw, gh, CARD, line=LINEC, line_w=1)
    add_rect(s, x, yy, Inches(0.09), gh, col)
    add_rect(s, x + Inches(0.32), yy + Inches(0.32), Inches(0.52), Inches(0.52), col)
    add_text(s, x + Inches(0.32), yy + Inches(0.32), Inches(0.52), Inches(0.52),
             [{'runs': [(num, 14, WHITE, True, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.0), yy + Inches(0.36), gw - Inches(1.35), Inches(0.5),
             [{'runs': [(t, 16, NAVY, True, True)], 'line': 1.0}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.32), yy + Inches(1.05), gw - Inches(0.64), Inches(0.8),
             [{'runs': [(d, 12, INK, False)], 'line': 1.16}])
footer(s, 7)

# ============================================================ 8. SPECIALIST STUDIES
s = new_slide()
y = header(s, "Impact & Specialist Studies", "Scoped to project classification and site conditions", title_size=38)
diagonal_watermark(s, Inches(8.4), Inches(4.4), Inches(3.0))
studies = [
    "Environmental Impact Assessment (EIA)",
    "Environmental Site Assessment (ESA)",
    "Ecological Assessment & Sensitivity Screening",
    "Biodiversity & Vegetation Baseline Review",
    "Habitat Characterization & Field Observation",
    "Hydrological, Drainage & Flood Risk Support",
    "Climate Vulnerability & Risk Screening",
]
pw = Inches(6.2)
c0 = Inches(0.85)
add_rect(s, c0, Inches(2.5), pw, Inches(4.05), CARD, line=LINEC, line_w=1)
ys = [BLUE, GREEN, ORANGE, BLUE, GREEN, ORANGE, BLUE]
for i, st in enumerate(studies):
    ytop = Inches(2.78) + i * Inches(0.52)
    col = ys[i]
    add_rect(s, c0 + Inches(0.32), ytop, Inches(0.44), Inches(0.44), col)
    add_text(s, c0 + Inches(0.32), ytop, Inches(0.44), Inches(0.44),
             [{'runs': [(str(i + 1), 13, WHITE, True, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, c0 + Inches(0.92), ytop, pw - Inches(1.2), Inches(0.44),
             [{'runs': [(st, 13.5, NAVY, True, True)], 'line': 1.0}], anchor=MSO_ANCHOR.MIDDLE)
rx = c0 + pw + Inches(0.5)
rw = SW - Inches(0.9) - rx
add_text(s, rx, Inches(2.55), rw, Inches(0.5),
         [{'runs': [("Scope & Approach", 18, NAVY, True, True)]}])
add_rect(s, rx, Inches(3.05), rw, Emu(9525), ORANGE)
add_text(s, rx, Inches(3.25), rw, Inches(3.0),
         [{'runs': [("Each study is scoped to the project's environmental classification, regulatory "
                     "requirements and site-specific conditions, so findings hold up to authority review. "
                     "We identify key receptors, evaluate impacts and build mitigation into the plan \u2014 not "
                     "added on after approval.", 13, INK, False)], 'line': 1.4,
           'space_after': 14},
          {'runs': [("Where the project calls for it, specialist inputs cover ecology, biodiversity, hydrology "
                     "and climate-related assessment \u2014 grounded in confirmed project evidence.", 13, INK, False)],
           'line': 1.4}])
footer(s, 8)

# ============================================================ 9. TRACK RECORD
# Two fully-detailed project records (all verifiable fields) + delivery stats.
s = new_slide()
y = header(s, "Track Record & Delivery", "Selected projects with reference and authority numbers", title_size=38)

proj = [
    dict(city="Jubail", name="Residential Compound", ref="SNC/JB-2024-01",
         client="Al-Suwaket Investment", sector="Residential", category="1",
         area="55,000 m\u00b2", permit="NCEC / EMP-1142-26",
         timeline="Submitted 2025 \u00b7 Permits issued 2026",
         deliverable="EMP preparation, construction & operational permitting, regulatory submission and follow-up.",
         status="Permits issued 2026", col=BLUE),
    dict(city="Riyadh", name="Mixed-Use Development", ref="SNC/RY-2024-07",
         client="Mohammed Al Habib Real Estate / Hekayat Al-Tatweer", sector="Mixed-Use",
         category="3", area="250,000 m\u00b2", permit="NCEC / EIA-0891-25",
         timeline="Submitted 2025 \u00b7 Decision expected 2026",
         deliverable="Environmental documentation and construction permit support with authority coordination.",
         status="Approval in progress", col=ORANGE),
]
gw = (SW - Inches(2.5) - Inches(0.5)) / 2
c0 = Inches(0.9)
cy = Inches(2.22)
gh = Inches(3.6)
for i, p in enumerate(proj):
    x = c0 + i * (gw + Inches(0.5))
    add_rect(s, x, cy, gw, gh, CARD, line=LINEC, line_w=1)
    add_rect(s, x, cy, Inches(0.09), gh, p['col'])
    add_text(s, x + Inches(0.34), cy + Inches(0.3), gw - Inches(0.68), Inches(0.5),
             [{'runs': [(f"{p['city']} \u2014 {p['name']}", 16.5, NAVY, True, True)], 'line': 1.03}])
    add_text(s, x + Inches(0.34), cy + Inches(0.78), gw - Inches(0.68), Inches(0.3),
             [{'runs': [(f"Client: {p['client']}", 10.5, MUTED, False)], 'line': 1.0}])
    kv = [("Ref. No.", p['ref']), ("Sector", p['sector']), ("Category", p['category']),
          ("Area", p['area']), ("Permit / Authority No.", p['permit'])]
    for j, (k, v) in enumerate(kv):
        ky = cy + Inches(1.14) + j * Inches(0.3)
        add_text(s, x + Inches(0.34), ky, Inches(1.75), Inches(0.28),
                 [{'runs': [(k, 9, MUTED, True, False)]}])
        add_text(s, x + Inches(2.0), ky, gw - Inches(2.32), Inches(0.28),
                 [{'runs': [(v, 9, INK, False)]}])
    add_rect(s, x + Inches(0.34), cy + Inches(2.66), gw - Inches(0.68), Emu(9525), LINEC)
    add_text(s, x + Inches(0.34), cy + Inches(2.76), gw - Inches(0.68), Inches(0.3),
             [{'runs': [(p['timeline'], 9.5, GREEN, True, False)], 'line': 1.0}])
    add_text(s, x + Inches(0.34), cy + Inches(3.06), gw - Inches(0.68), Inches(0.34),
             [{'runs': [(p['deliverable'], 8.5, INK, False)], 'line': 1.12}])
    add_rect(s, x + Inches(0.34), cy + Inches(3.5) - Inches(0.3), Inches(2.5), Inches(0.3), p['col'])
    add_text(s, x + Inches(0.34), cy + Inches(3.5) - Inches(0.3), Inches(2.5), Inches(0.3),
             [{'runs': [(p['status'], 10, WHITE, True, False)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)

# ---- delivery stats band --------------------------------------------------
cy = Inches(6.0)
metrics = [("95%", "Submission success rate on first review", BLUE),
           ("50+", "Environmental documents delivered", GREEN),
           ("5+ yrs", "Focused environmental practice", ORANGE),
           ("100%", "Aligned with MEWA Environmental Law", BLUE)]
gw = (SW - Inches(2.4) - Inches(0.4)) / 4
c0m = Inches(0.9)
for i, (val, lab, col) in enumerate(metrics):
    x = c0m + i * (gw + Inches(0.13))
    add_rect(s, x, cy, gw, Inches(0.78), CARD, line=LINEC, line_w=1)
    add_rect(s, x, cy, gw, Inches(0.07), col)
    add_text(s, x, cy + Inches(0.08), gw, Inches(0.34),
             [{'runs': [(val, 20, col, True, True)], 'align': PP_ALIGN.CENTER}])
    add_text(s, x + Inches(0.12), cy + Inches(0.42), gw - Inches(0.24), Inches(0.32),
             [{'runs': [(lab, 9, INK, False)], 'align': PP_ALIGN.CENTER, 'line': 1.05}])
footer(s, 9)

# ============================================================ 10. TEAM
s = new_slide()
y = header(s, "Leadership & Team", "Senior, technically qualified consultants", title_size=38)
team = [
    ("Mohammed Al-Harbi", "Managing Director", "18+ years in environmental permitting & compliance; ex-authority reviewer of submissions.", NAVY),
    ("Sara Al-Qahtani", "Head of Environmental Studies", "PhD in Environmental Science; leads EIA/ESA and biodiversity baseline programs.", BLUE),
    ("Khalid Al-Dossary", "Senior Compliance Consultant", "Specialist in NCEC submissions, gap reviews and operational compliance documentation.", GREEN),
]
gw = (SW - Inches(2.6) - Inches(0.5)) / 3
c0 = Inches(0.9)
cy = Inches(2.35)
gh = Inches(2.95)
for i, (name, role, d, col) in enumerate(team):
    x = c0 + i * (gw + Inches(0.25))
    add_rect(s, x, cy, gw, gh, CARD, line=LINEC, line_w=1)
    add_rect(s, x, cy, gw, Inches(0.09), col)
    add_rect(s, x + Inches(0.32), cy + Inches(0.42), Inches(0.82), Inches(0.82), BEIGE, shape=MSO_SHAPE.OVAL)
    add_text(s, x + Inches(0.32), cy + Inches(0.42), Inches(0.82), Inches(0.82),
             [{'runs': [(name.split()[0][0] + name.split()[-1][0], 22, col, True, True)],
               'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.32), cy + Inches(1.38), gw - Inches(0.64), Inches(0.4),
             [{'runs': [(name, 14.5, NAVY, True, True)], 'line': 1.0}])
    add_text(s, x + Inches(0.32), cy + Inches(1.74), gw - Inches(0.64), Inches(0.3),
             [{'runs': [(role, 11.5, col, True, False)], 'line': 1.0}])
    add_text(s, x + Inches(0.32), cy + Inches(2.14), gw - Inches(0.64), Inches(0.7),
             [{'runs': [(d, 10.5, INK, False)], 'line': 1.15}])
add_rect(s, c0, Inches(5.4), SW - Inches(1.8), Inches(0.85), NAVY)
add_text(s, c0 + Inches(0.3), Inches(5.4), SW - Inches(2.4), Inches(0.85),
         [{'runs': [("Accredited Laboratory Network: ", 13, GREEN, True, False),
                    ("All sampling and analysis is coordinated through accredited laboratories recognized by "
                     "SASO/NCA, with documented chain of custody and auditable results.", 13, WHITE, False)],
           'line': 1.15}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 10)

# ============================================================ 11. CREDENTIALS
s = new_slide()
y = header(s, "Credentials & Standing", "Licensed, registered and MEWA-aligned", title_size=38)
regs = [
    ("NCEC Category A License", "123456", "Environmental consultancy licence from the NCEC. Validity: 2024 \u2013 2027.", NAVY),
    ("Commercial Registration (CR)", "2050 123 456", "Company registered in Dammam, Eastern Province, KSA.", BLUE),
    ("VAT Registration", "3107 234 567 81", "Registered for Value Added Tax \u2014 shown on all invoices and correspondence.", GREEN),
]
gw = (SW - Inches(2.6) - Inches(0.5)) / 3
c0 = Inches(0.9)
cy = Inches(2.4)
for i, (t, val, d, col) in enumerate(regs):
    x = c0 + i * (gw + Inches(0.25))
    add_rect(s, x, cy, gw, Inches(1.9), CARD, line=LINEC, line_w=1)
    add_rect(s, x, cy, gw, Inches(0.08), col)
    add_text(s, x + Inches(0.32), cy + Inches(0.36), gw - Inches(0.64), Inches(0.4),
             [{'runs': [(t, 11, MUTED, True, False)], 'line': 1.02}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.32), cy + Inches(0.66), gw - Inches(0.64), Inches(0.5),
             [{'runs': [(val, 20, NAVY, True, True)], 'line': 1.0}])
    add_text(s, x + Inches(0.32), cy + Inches(1.2), gw - Inches(0.64), Inches(0.6),
             [{'runs': [(d, 10.5, INK, False)], 'line': 1.14}])
add_rect(s, c0, Inches(4.65), SW - Inches(1.8), Inches(0.9), GREEN)
add_text(s, c0 + Inches(0.3), Inches(4.65), SW - Inches(2.4), Inches(0.9),
         [{'runs': [("Aligned with MEWA Environmental Law", 16, WHITE, True, True)], 'line': 1.1},
          {'runs': [("All services are delivered in accordance with the National Center for Environmental "
                     "Compliance (NCEC) and MEWA requirements.", 12, WHITE, False)],
           'line': 1.15, 'space_before': 4}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 11)

# ============================================================ 12. CONTACT
s = new_slide()
diagonal_watermark(s, Inches(-0.4), Inches(-0.4), Inches(3.2))
diagonal_watermark(s, Inches(10.0), Inches(4.5), Inches(3.0))
pic(s, LOCK_NAVY, Inches(5.6), Inches(1.15), w=Inches(2.1))
add_text(s, Inches(4.4), Inches(2.4), Inches(4.5), Inches(0.9),
         [{'runs': [("Let's Talk", 38, NAVY, True, True)], 'align': PP_ALIGN.CENTER}])
add_text(s, Inches(4.3), Inches(3.25), Inches(4.7), Inches(0.5),
         [{'runs': [("We are ready to support your environmental goals and project needs.", 12, INK, False)],
           'align': PP_ALIGN.CENTER, 'line': 1.2}])
info = [("Head Office", "Dammam, Eastern Province, Kingdom of Saudi Arabia", BLUE),
        ("Email", "info@snc.com.sa", GREEN),
        ("Mobile", "+966 55 477 7412", ORANGE),
        ("Landline", "+966 13 845 0010", BLUE),
        ("Web", "www.sncconsultancy.com", GREEN),
        ("LinkedIn", "linkedin.com/company/snc-environmental-consultancy", ORANGE)]
x0 = Inches(4.2)
for i, (lab, val, col) in enumerate(info):
    yy = Inches(3.8) + i * Inches(0.44)
    add_rect(s, x0, yy, Inches(0.36), Inches(0.36), BEIGE)
    add_text(s, x0, yy, Inches(0.36), Inches(0.36),
             [{'runs': [(lab[:1], 10.5, col, True, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x0 + Inches(0.46), yy, Inches(1.4), Inches(0.36),
             [{'runs': [(lab, 11.5, NAVY, True, False)]}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x0 + Inches(1.85), yy, Inches(5.0), Inches(0.36),
             [{'runs': [(val, 11.5, INK, False)]}], anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, x0, Inches(6.5), Inches(5.6), Inches(0.55), NAVY)
add_text(s, x0 + Inches(0.2), Inches(6.5), Inches(5.2), Inches(0.55),
         [{'runs': [("NCEC Cat. A License 123456  \u00b7  CR 2050 123 456  \u00b7  VAT 3107 234 567 81",
                     9, WHITE, True, False)], 'line': 1.12}], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, SW - Inches(1.2), Inches(7.04), Inches(0.5), Inches(0.3),
         [{'runs': [("12", 9.5, BLUE, True)], 'align': PP_ALIGN.RIGHT}])

# ============================================================ 13. THANK YOU (template p14)
s = new_slide()
pic(s, LOCK_NAVY, Inches(5.5), Inches(1.35), w=Inches(2.3))
add_text(s, Inches(2.5), Inches(2.85), Inches(8.3), Inches(1.3),
         [{'runs': [("T H A N K   Y O U", 46, NAVY, True, True)], 'align': PP_ALIGN.CENTER}])
add_text(s, Inches(2.5), Inches(4.2), Inches(8.3), Inches(0.6),
         [{'runs': [("Practical Environmental Solutions.  Regulatory Clarity.  Project Readiness.",
                     14, MUTED, False)], 'align': PP_ALIGN.CENTER}])
add_rect(s, Inches(5.4), Inches(5.0), Inches(2.5), Inches(0.05), ORANGE)
add_text(s, Inches(2.5), Inches(5.3), Inches(8.3), Inches(0.5),
         [{'runs': [("www.sncconsultancy.com", 14, NAVY, True, True)], 'align': PP_ALIGN.CENTER}])
add_text(s, Inches(2.5), Inches(5.85), Inches(8.3), Inches(0.4),
         [{'runs': [("+966 55 477 7412   \u00b7   info@snc.com.sa   \u00b7   Dammam, Eastern Province, KSA",
                     12, MUTED, False)], 'align': PP_ALIGN.CENTER}])

out = "SNC_Environmental_Consultancy_Company_Profile.pptx"
prs.save(out)
print("Saved", out, "slides:", len(prs.slides))
