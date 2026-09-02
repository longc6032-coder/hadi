#!/usr/bin/env python3
"""Approximate visual renderer for the PPTX using PIL to catch layout/overflow issues."""
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw, ImageFont
import os

prs = Presentation('SNC_Environmental_Consultancy_Company_Profile.pptx')
SCALE = 120
W = int(prs.slide_width / 914400 * SCALE)
H = int(prs.slide_height / 914400 * SCALE)
os.makedirs('render_out', exist_ok=True)


def font(size, bold):
    base = '/usr/share/fonts/truetype/dejavu/DejaVuSans%s.ttf' % ('-Bold' if bold else '')
    if os.path.exists(base):
        return ImageFont.truetype(base, size)
    return ImageFont.load_default()


def e2px(v):
    return int(round(v / 914400 * SCALE))


def wrap(draw, text, font, w):
    words = text.split()
    lines, cur = [], ""
    for wd in words:
        t = (cur + " " + wd).strip()
        if draw.textlength(t, font=font) <= w:
            cur = t
        else:
            lines.append(cur); cur = wd
    lines.append(cur)
    return lines


def draw_text_block(draw, shape):
    x, y = e2px(shape.left), e2px(shape.top)
    w, h = e2px(shape.width), e2px(shape.height)
    tf = shape.text_frame
    anchor = tf.vertical_anchor
    paras = []
    cy_off = 0
    for para in tf.paragraphs:
        txt = "".join(r.text for r in para.runs)
        if not txt.strip():
            continue
        size, color, bold = 18, (36, 48, 58), False
        for r in para.runs:
            if r.font.size:
                size = r.font.size.pt
            if r.font.color and r.font.color.rgb:
                color = tuple(r.font.color.rgb[:3] if hasattr(r.font.color.rgb, '__getitem__') and len(r.font.color.rgb)>=3 else (36,48,58))
                try:
                    color = (r.font.color.rgb[0], r.font.color.rgb[1], r.font.color.rgb[2])
                except Exception:
                    color = (36,48,58)
            bold = bool(r.font.bold)
        f = font(int(size * SCALE / 72), bold)
        lh = int(size * SCALE / 72 * 1.18)
        lines = wrap(draw, txt, f, w)
        paras.append((lines, f, size, color, lh, para.alignment, para.space_before.pt if para.space_before else 0))
    # total height
    total = sum(len(l) * lh + sb for (l, f, size, color, lh, al, sb) in paras)
    cy = y
    if anchor == MSO_ANCHOR.MIDDLE:
        cy = y + max(0, (h - total) // 2)
    elif anchor == MSO_ANCHOR.BOTTOM:
        cy = y + max(0, (h - total))
    for (lines, f, size, color, lh, al, sb) in paras:
        cy += sb
        for ln in lines:
            tw = draw.textlength(ln, font=f)
            if al == PP_ALIGN.CENTER:
                tx = x + (w - tw) // 2
            elif al == PP_ALIGN.RIGHT:
                tx = x + (w - tw)
            else:
                tx = x
            draw.text((tx, cy), ln, fill=color, font=f)
            cy += lh
        cy += 6


for si, slide in enumerate(prs.slides):
    img = Image.new('RGB', (W, H), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    for shape in slide.shapes:
        if shape.shape_type == 13:  # picture
            try:
                import tempfile
                blob = shape.image.blob
                tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                tmp.write(blob); tmp.close()
                pic = Image.open(tmp.name).convert('RGBA')
                pic = pic.resize((e2px(shape.width), e2px(shape.height)))
                img.paste(pic, (e2px(shape.left), e2px(shape.top)), pic)
            except Exception as e:
                print('pict err', si + 1, e)
            continue
        left, top, w, h = e2px(shape.left), e2px(shape.top), e2px(shape.width), e2px(shape.height)
        fill = None
        line = None
        try:
            ft = str(shape.fill.type)
            if 'SOLID' in ft or str(shape.fill.type) == '1':
                fc = shape.fill.fore_color.rgb
                fill = (fc[0], fc[1], fc[2])
        except Exception:
            pass
        try:
            if str(shape.line.fill.type) == '1':
                lc = shape.line.color.rgb
                line = (lc[0], lc[1], lc[2])
        except Exception:
            pass
        # draw shape
        if fill or line:
            draw.rectangle([left, top, left + w, top + h], fill=fill, outline=line, width=1)
        elif shape.shape_type == 9:  # oval
            draw.ellipse([left, top, left + w, top + h], fill=fill, outline=line, width=1)
        if shape.has_text_frame and shape.text_frame.text.strip():
            draw_text_block(draw, shape)
    img.save(f'render_out/slide{si+1:02d}.png')
print('done')
